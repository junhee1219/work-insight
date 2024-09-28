from pptx import Presentation
import os


def ppt추출함수(파일경로, 추출단어, 제외슬라이드번호, 대소문자구분):
    파일명 = os.path.basename(파일경로).split('.')[0]
    
    prs = Presentation(파일경로)
    save_slide_ids = set()
    del_slides = set()
    for slide_number in 제외슬라이드번호:
        try:
            save_slide_ids.add(prs.slides[int(slide_number) - 1].slide_id)
        except IndexError:
            return
    if 대소문자구분 :
        for slide in prs.slides:
            breaker = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text and 추출단어 in shape.text:
                        save_slide_ids.add(slide.slide_id)
                        breaker = True
                        break
                if shape.has_table:
                    for row in shape.table.rows:
                        if breaker:
                            break
                        for cell in row.cells:
                            if 추출단어 in cell.text:
                                save_slide_ids.add(slide.slide_id)
                                breaker = True
                                break
    else:
        for slide in prs.slides:
            breaker = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text and 추출단어.upper() in shape.text.upper():
                        save_slide_ids.add(slide.slide_id)
                        breaker = True
                        break
                if shape.has_table:
                    for row in shape.table.rows:
                        if breaker:
                            break
                        for cell in row.cells:
                            if 추출단어.upper() in cell.text.upper():
                                save_slide_ids.add(slide.slide_id)
                                breaker = True
                                break
    for slide in prs.slides._sldIdLst:
        if not (int(slide.get("id")) in save_slide_ids):
            del_slides.add(slide)
    for i in del_slides:
        prs.slides._sldIdLst.remove(i)
    최종파일경로 = f"{파일명}_{추출단어}.pptx"
    prs.save(최종파일경로)
    return 최종파일경로


